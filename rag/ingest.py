"""Document ingestion and parsing."""
import os
import tempfile
from typing import List, Dict, Any
import fitz  # PyMuPDF
import docx  # python-docx

try:
    from pptx import Presentation  # python-pptx
except Exception:
    Presentation = None

try:
    import win32com.client  # pywin32 (Windows only), for converting .ppt -> .pptx
except Exception:
    win32com = None


class DocumentParser:
    """Parse documents for RAG."""
    
    @staticmethod
    def parse_pdf(file_path: str) -> List[Dict[str, Any]]:
        """Parse PDF and extract text with page numbers."""
        pages = []
        try:
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text()
                if text.strip():
                    pages.append({
                        "text": text,
                        "page": page_num + 1,
                        "doc_id": os.path.basename(file_path)
                    })
            doc.close()
        except Exception as e:
            print(f"Error parsing PDF {file_path}: {e}")
        return pages
    
    @staticmethod
    def parse_txt(file_path: str) -> List[Dict[str, Any]]:
        """Parse text file with encoding auto-detection (UTF-8 → GBK → Latin-1)."""
        text = None
        for encoding in ('utf-8-sig', 'utf-8', 'gbk', 'latin-1'):
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    text = f.read()
                break
            except UnicodeDecodeError:
                continue
            except Exception as e:
                print(f"Error parsing TXT {file_path}: {e}")
                return []
        if text is None:
            print(f"Error parsing TXT {file_path}: 无法识别文件编码")
            return []
        return [{
            "text": text,
            "page": None,
            "doc_id": os.path.basename(file_path)
        }]
    
    @staticmethod
    def parse_docx(file_path: str) -> List[Dict[str, Any]]:
        """Parse Word document (.docx)."""
        try:
            doc = docx.Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text = "\n".join(paragraphs)
            return [{
                "text": text,
                "page": None,
                "doc_id": os.path.basename(file_path)
            }]
        except Exception as e:
            print(f"Error parsing DOCX {file_path}: {e}")
            return []

    @staticmethod
    def parse_pptx(file_path: str) -> List[Dict[str, Any]]:
        """Parse PowerPoint document (.pptx). Each slide is treated as one page."""
        if Presentation is None:
            print("Error parsing PPTX: python-pptx 未安装")
            return []

        try:
            prs = Presentation(file_path)
            pages = []
            for slide_idx, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            texts.append(text)

                slide_text = "\n".join(texts).strip()
                if slide_text:
                    pages.append({
                        "text": slide_text,
                        "page": slide_idx,
                        "doc_id": os.path.basename(file_path)
                    })
            return pages
        except Exception as e:
            print(f"Error parsing PPTX {file_path}: {e}")
            return []

    @staticmethod
    def _convert_ppt_to_pptx(file_path: str) -> str:
        """Convert legacy .ppt to .pptx via PowerPoint COM (Windows). Returns temp .pptx path."""
        if win32com is None:
            raise RuntimeError("pywin32 未安装，无法解析 .ppt，请安装 pywin32 或先手动转为 .pptx")

        fd, temp_pptx = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)

        app = None
        presentation = None
        try:
            app = win32com.client.Dispatch("PowerPoint.Application")
            presentation = app.Presentations.Open(
                file_path,
                ReadOnly=1,
                Untitled=0,
                WithWindow=0,
            )
            # 24 = ppSaveAsOpenXMLPresentation
            presentation.SaveAs(temp_pptx, 24)
            return temp_pptx
        finally:
            try:
                if presentation is not None:
                    presentation.Close()
            except Exception:
                pass
            try:
                if app is not None:
                    app.Quit()
            except Exception:
                pass

    @staticmethod
    def parse_ppt(file_path: str) -> List[Dict[str, Any]]:
        """Parse legacy .ppt by converting to .pptx first."""
        temp_pptx = None
        try:
            temp_pptx = DocumentParser._convert_ppt_to_pptx(file_path)
            return DocumentParser.parse_pptx(temp_pptx)
        except Exception as e:
            print(f"Error parsing PPT {file_path}: {e}")
            return []
        finally:
            if temp_pptx and os.path.exists(temp_pptx):
                try:
                    os.remove(temp_pptx)
                except Exception:
                    pass

    @staticmethod
    def parse_document(file_path: str) -> List[Dict[str, Any]]:
        """Parse document based on file extension."""
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.pdf':
            return DocumentParser.parse_pdf(file_path)
        elif ext in ['.txt', '.md']:
            return DocumentParser.parse_txt(file_path)
        elif ext == '.docx':
            return DocumentParser.parse_docx(file_path)
        elif ext == '.pptx':
            return DocumentParser.parse_pptx(file_path)
        elif ext == '.ppt':
            return DocumentParser.parse_ppt(file_path)
        else:
            print(f"Unsupported file type: {ext}")
            return []
